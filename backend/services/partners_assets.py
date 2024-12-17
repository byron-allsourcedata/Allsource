import logging
import os
from typing import Union
import tempfile
import zipfile
from pathlib import Path
from PIL import Image, ImageDraw, UnidentifiedImageError
from pptx import Presentation
import subprocess
from fastapi import HTTPException, UploadFile
from io import BytesIO
from urllib.parse import urlparse
import requests
from enums import PartnersAssetsInfoEnum
from models.partners_asset import PartnersAsset
from services.aws import AWSService
from persistence.partners_asset_persistence import PartnersAssetPersistence
from schemas.partners_asset import PartnersAssetResponse

logger = logging.getLogger(__name__)


class PartnersAssetService:

    def __init__(self,
        partners_asset_persistence: PartnersAssetPersistence,
        aws_service: AWSService):
        self.partners_asset_persistence = partners_asset_persistence
        self.AWS = aws_service


    def get_assets(self):
        assets = self.partners_asset_persistence.get_assets()
        return [
            self.domain_mapped(asset)
            for i, asset in enumerate(assets)
        ]
    

    def delete_asset(self, id: int):
        if id:
            try:
                self.partners_asset_persistence.delete_asset(asset_id=id)
                return PartnersAssetsInfoEnum.SUCCESS
            except Exception as err:
                logger.debug('Error deleting assets file', err)
                return PartnersAssetsInfoEnum.NOT_FOUND
        else:
            return PartnersAssetsInfoEnum.NOT_VALID_ID


    async def update_asset(self, asset_id: int, description: str, file: UploadFile = None):
        updating_data = {"description": description}

        if file:
            file_contents = await file.read()
            file_extension = Path(file.filename).suffix.lower()
            self.AWS.upload_string(file_contents, f'partners-assets/{description}{file_extension}')
            preview = await self.generate_preview(file_contents, file_extension, type)
            if preview:
                file_preview_contents = preview.read()
                self.AWS.upload_string(file_preview_contents, f'partners-assets/{description}_preview{file_extension}')
            updating_data["file_url"] = f'https://maximiz-data.s3.us-east-2.amazonaws.com/partners-assets/{description}{file_extension}'
            updating_data["preview_url"] = f'https://maximiz-data.s3.us-east-2.amazonaws.com/partners-assets/{description}_preview{file_extension}'
 
        updated_data = self.partners_asset_persistence.update_asset(asset_id, updating_data)

        if not updated_data:
            raise HTTPException(status_code=500, detail="Asset not updated")

        return self.domain_mapped(updated_data)
    

    async def upload_files_asset(self, description: str, file: UploadFile):
        file_contents = await file.read()
        file_extension = Path(file.filename).suffix.lower()
        self.AWS.upload_string(file_contents, f'partners-assets/{description}{file_extension}')
        preview = await self.generate_preview(file_contents, file_extension, type)
        if preview:
            file_preview_contents = preview.read()
            self.AWS.upload_string(file_preview_contents, f'partners-assets/{description}_preview{file_extension}')


    async def create_asset(self, description: str, type: str, file: UploadFile = None):
        if(file):
            file_contents = await file.read()
            file_extension = Path(file.filename).suffix.lower()
            self.AWS.upload_string(file_contents, f'partners-assets/{description}{file_extension}')
            preview = await self.generate_preview(file_contents, file_extension, type)
            if preview:
                file_preview_contents = preview.read()
                self.AWS.upload_string(file_preview_contents, f'partners-assets/{description}_preview{file_extension}')
            creating_data = {
                "description": description, 
                "type": type, 
                "file_url": f'https://maximiz-data.s3.us-east-2.amazonaws.com/partners-assets/{description}{file_extension}',
                "preview_url": f'https://maximiz-data.s3.us-east-2.amazonaws.com/partners-assets/{description}_preview{file_extension}',
            }

            created_data = self.partners_asset_persistence.create_data(creating_data)

            if not created_data:
                raise HTTPException(status_code=500, detail="Asset not created")
            
            return self.domain_mapped(created_data)
        else: 
            raise HTTPException(status_code=404, detail="File not received")

    
    async def generate_preview(self, file_contents: bytes, file_extension: str, type_asset: str) -> Union[BytesIO, None]:
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            temp_file.write(file_contents)
            temp_file_path = temp_file.name

        if type_asset == "image":
            return self._generate_image_preview(temp_file_path)

        elif type_asset == "video":
            return self._generate_video_preview(temp_file_path)

        elif type_asset == "presentation":
            return self._generate_presentation_preview(temp_file_path)
        else:
            return None
    

    def _generate_image_preview(self, file_path: str) -> BytesIO:
        try:
            with Image.open(file_path) as img:
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.thumbnail((300, 300))
                preview = BytesIO()
                img.save(preview, format="JPEG")
                preview.seek(0)
                return preview
        except UnidentifiedImageError:
            raise HTTPException(status_code=400, detail="Invalid image file")


    def _generate_video_preview(self, file_path: str) -> BytesIO:
        preview_path = f"{file_path}_preview.jpg"
        subprocess.run(
            [
                "ffmpeg",
                "-i", file_path,
                "-ss", "00:00:01.000",
                "-vframes", "1",
                preview_path,
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )

        with open(preview_path, "rb") as preview_file:
            preview = BytesIO(preview_file.read())
        return preview
    

    # def _generate_presentation_preview(self, file_path: str) -> BytesIO:
    #     prs = Presentation(file_path)
    #     first_slide = prs.slides[0]

    #     preview_image = None

    #     background_image = self._extract_slide_background(file_path, slide_index=0)
    #     if background_image:
    #         preview_image = Image.open(background_image)
    #         preview_image.thumbnail((1920, 1080))

  
    #     if preview_image is None:
    #         for shape in first_slide.shapes:
    #             if shape.shape_type == 13: 
    #                 img_stream = BytesIO(shape.image.blob)
    #                 img = Image.open(img_stream)
    #                 img.thumbnail((1920, 1080))
    #                 preview_image = img
    #                 break

    #     if preview_image is None:
    #         preview_image = Image.new("RGB", (1920, 1080), (255, 255, 255))
    #         draw = ImageDraw.Draw(preview_image)
    #         draw.text((960, 540), "No Preview Available", fill="black", anchor="mm")


    #     preview_image.thumbnail((300, 300))
    #     preview = BytesIO()
    #     preview_image.save(preview, format="JPEG")
    #     preview.seek(0)
    #     return preview


    # def _extract_slide_background(self, file_path: str, slide_index: int = 0) -> Union[BytesIO, None]:
    #     with zipfile.ZipFile(file_path, 'r') as archive:
    #         slide_path = f'ppt/slides/slide{slide_index + 1}.xml'
    #         if slide_path not in archive.namelist():
    #             return None


    #         rels_path = f'ppt/slides/_rels/slide{slide_index + 1}.xml.rels'
    #         if rels_path in archive.namelist():
    #             rels_content = archive.read(rels_path).decode()
    #             if 'image' in rels_content:
    #                 media_path = rels_content.split('Target="..')[1].split('"')[0]
    #                 media_full_path = os.path.join('ppt', media_path.strip('/'))

    #                 if media_full_path in archive.namelist():
    #                     with archive.open(media_full_path) as media_file:
    #                         return BytesIO(media_file.read())
    #     return None


    # def _generate_presentation_preview(self, file_path: str) -> BytesIO:
    #     prs = Presentation(file_path)
    #     first_slide = prs.slides[0]
    #     # image = Image.new("RGB", (1920, 1080), (255, 255, 255))
    #     # preview = BytesIO()
    #     # image.thumbnail((300, 300))
    #     # image.save(preview, format="JPEG")
    #     # preview.seek(0)
    #     # return preview

    #     preview_image = None

    #     if first_slide.background and first_slide.background.fill.type == 1:
    #         color = first_slide.background.fill.fore_color.rgb
    #         preview_image = Image.new("RGB", (1920, 1080), (color.red, color.green, color.blue))
        
    #     preview_image.thumbnail((300, 300))
    #     preview = BytesIO()
    #     preview_image.save(preview, format="JPEG")
    #     preview.seek(0)
    #     return preview
    
        # images = []
        # for shape in first_slide.shapes:
        #     if shape.shape_type == 13: 
        #         images.append(shape.image)

        # if not images:
        #     raise ValueError("No images found on the first slide.")

        # image_stream = BytesIO(images[0].blob)
        # img = Image.open(image_stream)

        # img.thumbnail((300, 300))
        # preview = BytesIO()
        # img.save(preview, format="JPEG")
        # preview.seek(0)
        # return preview


    def get_file_size(self, file_url: str) -> str:
        if file_url:
            try:
                response = requests.head(file_url, allow_redirects=True, timeout=5)
                if response.status_code == 200 and 'Content-Length' in response.headers:
                    size_in_bytes = int(response.headers['Content-Length'])
                    size_in_mb = size_in_bytes / (1024 ** 2)
                    return f"{size_in_mb:.2f} MB"
                else:
                    return "0.00 MB"
            except Exception as err:
                logger.debug('Error fetching file size', err)
                return "0.00 MB" 
        else:
            return "0.00 MB"


    def get_file_extension(self, file_url) -> str:
        if file_url:
            url_path = urlparse(file_url).path
            extension = os.path.splitext(url_path)[-1][1:].capitalize()
            if not extension:
                return "Unknown"
            return extension
        else: 
            return "Unknown"


    def domain_mapped(self, asset: PartnersAsset):
        return PartnersAssetResponse(
            id=asset.id,
            title=asset.title,
            type=asset.type,
            preview_url=asset.preview_url,
            file_url=asset.file_url,
            file_extension=self.get_file_extension(asset.file_url),
            file_size=self.get_file_size(asset.file_url),
        ).model_dump()